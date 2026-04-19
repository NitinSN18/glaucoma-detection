"""
Generate a comprehensive 10-slide PowerPoint presentation about the Glaucoma Detection Project.
Covers project overview, models, Python libraries, code architecture, and results.
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
import os

def create_presentation():
    """Create a 10-slide PowerPoint presentation."""
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    
    # Define color scheme
    PRIMARY_COLOR = RGBColor(91, 216, 232)  # Cyan
    DARK_BG = RGBColor(20, 30, 45)
    TEXT_PRIMARY = RGBColor(255, 255, 255)
    TEXT_SECONDARY = RGBColor(200, 200, 200)
    ACCENT_GREEN = RGBColor(76, 175, 80)
    ACCENT_ORANGE = RGBColor(255, 152, 0)
    
    def add_title_slide(title, subtitle):
        """Add a title slide."""
        slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = DARK_BG
        
        # Add title
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(9), Inches(1.5))
        title_frame = title_box.text_frame
        title_frame.text = title
        title_para = title_frame.paragraphs[0]
        title_para.font.size = Pt(54)
        title_para.font.bold = True
        title_para.font.color.rgb = PRIMARY_COLOR
        title_para.alignment = PP_ALIGN.CENTER
        
        # Add subtitle
        subtitle_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.2), Inches(9), Inches(2))
        subtitle_frame = subtitle_box.text_frame
        subtitle_frame.word_wrap = True
        subtitle_frame.text = subtitle
        subtitle_para = subtitle_frame.paragraphs[0]
        subtitle_para.font.size = Pt(24)
        subtitle_para.font.color.rgb = TEXT_SECONDARY
        subtitle_para.alignment = PP_ALIGN.CENTER
        
        return slide
    
    def add_content_slide(title, content_items, bg_color=DARK_BG):
        """Add a content slide with bullet points."""
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = bg_color
        
        # Add title
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(9), Inches(0.8))
        title_frame = title_box.text_frame
        title_frame.text = title
        title_para = title_frame.paragraphs[0]
        title_para.font.size = Pt(40)
        title_para.font.bold = True
        title_para.font.color.rgb = PRIMARY_COLOR
        
        # Add horizontal line
        line = slide.shapes.add_shape(1, Inches(0.5), Inches(1.3), Inches(9), Inches(0))
        line.line.color.rgb = PRIMARY_COLOR
        line.line.width = Pt(2)
        
        # Add content
        content_box = slide.shapes.add_textbox(Inches(0.7), Inches(1.6), Inches(8.6), Inches(5.5))
        text_frame = content_box.text_frame
        text_frame.word_wrap = True
        
        for i, item in enumerate(content_items):
            if i > 0:
                text_frame.add_paragraph()
            p = text_frame.paragraphs[i]
            p.text = item
            p.font.size = Pt(18)
            p.font.color.rgb = TEXT_PRIMARY
            p.space_before = Pt(8)
            p.space_after = Pt(8)
            p.level = 0
        
        return slide
    
    # Slide 1: Title Slide
    add_title_slide(
        "Glaucoma Detection System",
        "AI-Powered Retinal Analysis using Deep Learning"
    )
    
    # Slide 2: Project Overview
    add_content_slide(
        "Project Overview",
        [
            "• Objective: Automated glaucoma screening from retinal fundus images",
            "• Impact: Early detection reduces vision loss by 50-80%",
            "• Approach: Dual AI pipeline (Classification + Segmentation)",
            "• Classification: EfficientNet-B0 for glaucoma risk prediction",
            "• Segmentation: DeepLabV3+ to extract optic disc and cup regions",
            "• Integration: Multi-pass refinement for robust mask generation"
        ]
    )
    
    # Slide 3: Clinical Problem & Solution
    add_content_slide(
        "Clinical Problem & Solution",
        [
            "Problem:",
            "  Glaucoma is asymptomatic until advanced (irreversible damage)",
            "  Manual screening is time-consuming and operator-dependent",
            "  Cup-to-Disc Ratio (CDR) is a key glaucoma indicator",
            "",
            "Solution:",
            "  Automated segmentation of optic disc and cup structures",
            "  Machine learning classifier trained on labeled fundus images",
            "  Real-time analysis with interpretable anatomical overlays"
        ]
    )
    
    # Slide 4: Models Architecture
    add_content_slide(
        "Deep Learning Models",
        [
            "Classification Model: EfficientNet-B0",
            "  • Efficient architecture (4.0M parameters) optimized for mobile/edge",
            "  • Pre-trained on ImageNet, fine-tuned on glaucoma fundus dataset",
            "  • Output: Binary classification (Glaucoma / Normal) with probability",
            "",
            "Segmentation Model: DeepLabV3+",
            "  • ResNet34 encoder with atrous spatial pyramid pooling (ASPP)",
            "  • 3-class output: Background, Optic Disc, Optic Cup",
            "  • Atrous convolution captures multi-scale anatomical features"
        ]
    )
    
    # Slide 5: Python Libraries & Dependencies
    add_content_slide(
        "Python Libraries & Tech Stack",
        [
            "Deep Learning Framework:",
            "  • PyTorch: Model inference, tensor operations, GPU/MPS acceleration",
            "",
            "Image Processing & Analysis:",
            "  • OpenCV (cv2): Morphological operations, contour detection, CLAHE enhancement",
            "  • Pillow (PIL): Image loading, resizing, format conversion",
            "",
            "Web Framework:",
            "  • Flask: REST API backend, session management, multi-endpoint routing",
            "  • Werkzeug: Secure file uploads, request/response handling",
            "",
            "Data & Utilities:",
            "  • NumPy: Numerical computations, matrix operations, probability fusion",
            "  • JSON: Patient record serialization, API responses"
        ]
    )
    
    # Slide 6: Segmentation Pipeline Architecture
    add_content_slide(
        "Segmentation Pipeline (5-Stage Fallback)",
        [
            "Stage 1: DeepLabV3+ Model Inference",
            "  → Multi-scale prediction + probability fusion across 3 crops",
            "",
            "Stage 2: Intensity-Based Refinement",
            "  → Extract disc/cup from brightness maps when model masks collapse",
            "",
            "Stage 3: Morphological Operations",
            "  → Erosion, dilation, connected component filtering for shape regularization",
            "",
            "Stage 4: Interior Stabilization",
            "  → Ensure cup stays within disc, smooth interior boundaries",
            "",
            "Stage 5: Ellipse Regularization",
            "  → Fit best-fit ellipses to enforce circular anatomical structure"
        ]
    )
    
    # Slide 7: Classification Pipeline & Feature Extraction
    add_content_slide(
        "Classification Pipeline",
        [
            "Input Preprocessing:",
            "  • Resize fundus image to 224×224 pixels (EfficientNet input size)",
            "  • Normalize with ImageNet statistics (mean=[0.485, 0.456, 0.406])",
            "",
            "Feature Extraction (EfficientNet-B0):",
            "  • 16 convolutional blocks → hierarchical feature learning",
            "  • Early layers: Edge detection, vessel patterns",
            "  • Mid layers: Cup/disc shape, texture features",
            "  • Final layers: Disease-specific patterns",
            "",
            "Classification Output:",
            "  • Sigmoid activation → Binary probability (0-1 scale)",
            "  • Softmax alternative for multi-class (Normal/Suspect/Glaucoma)"
        ]
    )
    
    # Slide 8: Code Structure & Key Functions
    add_content_slide(
        "Core Code Architecture (app.py)",
        [
            "Model Loading:",
            "  • load_classification_model(): Lazy-load EfficientNet-B0 on startup",
            "  • load_segmentation_model(): Load DeepLabV3+ with ResNet34 encoder",
            "",
            "Inference Functions:",
            "  • run_segmentation_pipeline(): Multi-pass mask refinement with fallback chain",
            "  • refine_disc_from_intensity(): Intensity-based recovery for collapsed masks",
            "  • choose_component_near_point(): Select anatomical structures near anchors",
            "",
            "API Endpoints:",
            "  • /api/combined/classify/segment: Full dual-pipeline analysis",
            "  • /api/classify: Classification-only mode",
            "  • /api/segment: Segmentation-only mode with quality scoring"
        ]
    )
    
    # Slide 9: Results Visualization & Quality Metrics
    add_content_slide(
        "Metrics & Interpretability",
        [
            "Classification Output:",
            "  • Glaucoma probability (0-100%)",
            "  • Sensitivity, Specificity, ROC-AUC scores",
            "  • Recommendation: 'Normal', 'Suspected Glaucoma', 'High Risk'",
            "",
            "Segmentation Output:",
            "  • Optic Disc Mask (green overlay on retina)",
            "  • Optic Cup Mask (inner yellow region)",
            "  • Cup-to-Disc Ratio (CDR): Clinical indicator of glaucoma",
            "  • Segmentation Quality Score (0-1): Model confidence metric",
            "",
            "Anatomical Anchors:",
            "  • Brightness prior: Gaussian centered on brightest region (optic disc location)",
            "  • Contour analysis: Circularity, solidity, area ratios for validation"
        ]
    )
    
    # Slide 10: Deployment & Future Enhancements
    add_content_slide(
        "Deployment & Roadmap",
        [
            "Current Deployment:",
            "  • Flask REST API with multi-device support (CPU, GPU, MPS)",
            "  • Session-based authentication for clinic workflows",
            "  • Patient record logging (JSONL format) for audit trails",
            "",
            "Integration Points:",
            "  • Single-image analysis for rapid screening",
            "  • Batch processing for high-throughput studies",
            "  • Patient logbook for longitudinal tracking",
            "",
            "Future Enhancements:",
            "  • Multi-modal fusion: Combine 2D segmentation with OCT/3D data",
            "  • Transfer learning: Fine-tune on hospital-specific demographics",
            "  • Model ensembling: Combine multiple architectures for higher accuracy",
            "  • Mobile deployment: ONNX export for on-device inference"
        ]
    )
    
    return prs

def save_presentation(output_path='glaucoma_presentation.pptx'):
    """Generate and save the presentation."""
    prs = create_presentation()
    prs.save(output_path)
    print(f"✓ Presentation saved to: {output_path}")
    return output_path

if __name__ == '__main__':
    save_presentation()
